"""填充 PPT 剩余空白页（5=过渡页导航, 11=总结与展望内容, 封面署名）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

PPT = r"C:\Users\ROG\Desktop\依赖安全与合规分析平台实习成果汇报.pptx"
OUT = r"C:\Users\ROG\Desktop\依赖安全与合规分析平台实习成果汇报_filled.pptx"
p = Presentation(PPT)

# ========== 幻灯片 1：封面改汇报人 ==========
s1 = p.slides[0]
for sh in s1.shapes:
    if sh.has_text_frame and "Alex" in sh.text_frame.text:
        for p2 in sh.text_frame.paragraphs:
            for run in p2.runs:
                if "Alex" in run.text:
                    run.text = run.text.replace("Alex", "张伟")
        break

# ========== 幻灯片 5（index 4）：过渡页加导航标签 ==========
s5 = p.slides[4]
# 在右侧添加两个白色文字标签："开发时间线回顾" 和 "核心指标速览"
# 位置参考 slide 3 的标签 (5.5,3.8) 和 (7.8,3.8)
nav_items = [("开发时间线回顾", 5.5), ("核心指标速览", 8.2)]
for text, left in nav_items:
    box = s5.shapes.add_textbox(Inches(left), Inches(3.8), Inches(2.3), Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = text
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

# ========== 幻灯片 11（index 10）：总结与展望内容页 ==========
s11 = p.slides[10]
# 已有：标题在 (5.5,2.7)，两个导航标签 (5.5,3.8) "项目总结" 和 (8.3,3.8) "个人收获与展望"
# 需要加实际内容

# 左栏："项目总结"
box1 = s11.shapes.add_textbox(Inches(4.0), Inches(4.4), Inches(5.0), Inches(2.5))
tf1 = box1.text_frame
tf1.word_wrap = True
lines_left = [
    "✅ 交付 v1.0.0 正式版：前后端完整平台，67/67 测试通过",
    "✅ 独立完成全栈：需求→架构→开发→测试→审计→发布",
    "✅ 把安全检查门槛从「安全专家」降到「5 分钟上手」",
    "✅ 8 周完成原定 12 周的核心开发，提前交付",
]
for i, line in enumerate(lines_left):
    if i == 0:
        p2 = tf1.paragraphs[0]
    else:
        p2 = tf1.add_paragraph()
    p2.text = line
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2.space_after = Pt(6)

# 右栏："收获与展望"
box2 = s11.shapes.add_textbox(Inches(4.0), Inches(4.4), Inches(5.0), Inches(2.5))
tf2 = box2.text_frame
tf2.word_wrap = True

p2 = tf2.paragraphs[0]
p2.text = "🎯 工程能力"
p2.font.size = Pt(12)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

p2 = tf2.add_paragraph()
p2.text = "Spring 生态 + 前端工程化 + 安全防护 + 全链路测试"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p2.space_after = Pt(8)

p2 = tf2.add_paragraph()
p2.text = "🧠 攻坚方法论"
p2.font.size = Pt(12)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

p2 = tf2.add_paragraph()
p2.text = "现象→假设→验证→绕行，NVD 离线方案为典型案例"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p2.space_after = Pt(8)

p2 = tf2.add_paragraph()
p2.text = "📋 文档与自审习惯"
p2.font.size = Pt(12)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

p2 = tf2.add_paragraph()
p2.text = "199 项开发核查点逐项审计，46 个问题闭环修复"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

p.save(OUT)
print(f"DONE — PPT 已保存到: {OUT}")
